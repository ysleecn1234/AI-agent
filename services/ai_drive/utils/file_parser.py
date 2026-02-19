"""
AI 드라이브 - 파일 파싱 유틸리티
지원 형식: PDF, DOCX, PPT, TXT, MD, CSV
"""

import pdfplumber  # 한글 PDF 파싱 지원
try:
    import fitz  # PyMuPDF (fallback용)
    HAS_FITZ = True
except ImportError:
    HAS_FITZ = False
from docx import Document
from pptx import Presentation
import csv
from pathlib import Path
from openpyxl import load_workbook

class FileParser:
    """
    다양한 파일 형식에서 텍스트 추출
    """
    
    def __init__(self):
        self.supported_formats = ['.pdf', '.docx', '.pptx', '.txt', '.md', '.csv', '.xlsx']
    
    def parse(self, file_path: str) -> str:
        """
        파일 형식에 따라 텍스트 추출
        
        Args:
            file_path: 파일 경로
            
        Returns:
            추출된 텍스트
            
        Raises:
            ValueError: 지원하지 않는 파일 형식
        """
        path = Path(file_path)
        suffix = path.suffix.lower()
        
        if suffix == '.pdf':
            return self.parse_pdf(file_path)
        elif suffix == '.docx':
            return self.parse_docx(file_path)
        elif suffix == '.pptx':
            return self.parse_pptx(file_path)
        elif suffix == '.txt':
            return self.parse_txt(file_path)
        elif suffix == '.md':
            return self.parse_md(file_path)
        elif suffix == '.csv':
            return self.parse_csv(file_path)
        elif suffix == '.xlsx':
            return self.parse_xlsx(file_path)
        else:
            raise ValueError(f"지원하지 않는 파일 형식: {suffix}")
    
    def parse_pdf(self, file_path: str) -> str:
        """
        PDF 파일에서 텍스트 추출 (pdfplumber 사용, 한글 지원)
        
        Args:
            file_path: PDF 파일 경로
            
        Returns:
            추출된 텍스트
        """
        text_content = []
        
        try:
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    text = page.extract_text()
                    
                    if text and text.strip():
                        text_content.append(f"[Page {page_num + 1}]\n{text}")
            
        except Exception as e:
            # pdfplumber 실패 시 fitz fallback
            if HAS_FITZ:
                print(f"  ⚠️ pdfplumber 실패, PyMuPDF로 재시도: {e}")
                return self._parse_pdf_fitz(file_path)
            raise Exception(f"PDF 파싱 오류: {str(e)}")
        
        return "\n\n".join(text_content)
    
    def _parse_pdf_fitz(self, file_path: str) -> str:
        """PyMuPDF fallback"""
        text_content = []
        doc = fitz.open(file_path)
        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            if text.strip():
                text_content.append(f"[Page {page_num + 1}]\n{text}")
        doc.close()
        return "\n\n".join(text_content)
    
    def parse_docx(self, file_path: str) -> str:
        """
        DOCX 파일에서 텍스트 추출
        
        Args:
            file_path: DOCX 파일 경로
            
        Returns:
            추출된 텍스트
        """
        try:
            doc = Document(file_path)
            
            text_content = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    text_content.append(para.text)
            
            return "\n\n".join(text_content)
            
        except Exception as e:
            raise Exception(f"DOCX 파싱 오류: {str(e)}")
    
    def parse_pptx(self, file_path: str) -> str:
        """
        PPTX 파일에서 텍스트 추출
        
        Args:
            file_path: PPTX 파일 경로
            
        Returns:
            추출된 텍스트
        """
        try:
            prs = Presentation(file_path)
            
            text_content = []
            
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_text = [f"[Slide {slide_num}]"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if len(slide_text) > 1:  # 제목 외에 내용이 있는 경우만
                    text_content.append("\n".join(slide_text))
            
            return "\n\n".join(text_content)
            
        except Exception as e:
            raise Exception(f"PPTX 파싱 오류: {str(e)}")
    
    def parse_txt(self, file_path: str) -> str:
        """
        TXT 파일에서 텍스트 추출
        
        Args:
            file_path: TXT 파일 경로
            
        Returns:
            추출된 텍스트
        """
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except UnicodeDecodeError:
            # UTF-8 실패 시 다른 인코딩 시도
            try:
                with open(file_path, 'r', encoding='cp949') as f:
                    return f.read()
            except Exception as e:
                raise Exception(f"TXT 파싱 오류: {str(e)}")
    
    def parse_md(self, file_path: str) -> str:
        """
        Markdown 파일에서 텍스트 추출
        
        Args:
            file_path: MD 파일 경로
            
        Returns:
            추출된 텍스트
        """
        # Markdown도 텍스트 파일이므로 txt와 동일하게 처리
        return self.parse_txt(file_path)
    
    def parse_csv(self, file_path: str) -> str:
        """
        CSV 파일에서 텍스트 추출
        
        Args:
            file_path: CSV 파일 경로
            
        Returns:
            추출된 텍스트 (테이블 형식)
        """
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                reader = csv.reader(f)
                rows = list(reader)
                
                if not rows:
                    return ""
                
                # 헤더와 데이터 구분
                text_content = []
                
                # 헤더
                if rows:
                    header = " | ".join(rows[0])
                    text_content.append(f"[CSV Header]\n{header}")
                
                # 데이터 행
                for idx, row in enumerate(rows[1:], start=1):
                    row_text = " | ".join(row)
                    text_content.append(f"[Row {idx}] {row_text}")
                
                return "\n".join(text_content)
                
        except Exception as e:
            raise Exception(f"CSV 파싱 오류: {str(e)}")
    
    def parse_xlsx(self, file_path: str) -> str:
        """
        Excel 파일에서 텍스트 추출
        
        Args:
            file_path: XLSX 파일 경로
            
        Returns:
            추출된 텍스트
        """
        try:
            wb = load_workbook(file_path, data_only=True)
            text_content = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_text = [f"[Sheet: {sheet_name}]"]
                
                for row in sheet.iter_rows():
                    row_values = []
                    for cell in row:
                        if cell.value is not None:
                            row_values.append(str(cell.value))
                    
                    if row_values:
                        sheet_text.append(" | ".join(row_values))
                
                if len(sheet_text) > 1:
                    text_content.append("\n".join(sheet_text))
            
            wb.close()
            return "\n\n".join(text_content)
            
        except Exception as e:
            raise Exception(f"XLSX 파싱 오류: {str(e)}")
    
    def is_supported(self, file_path: str) -> bool:
        """
        지원하는 파일 형식인지 확인
        
        Args:
            file_path: 파일 경로
            
        Returns:
            지원 여부
        """
        suffix = Path(file_path).suffix.lower()
        return suffix in self.supported_formats


# 테스트 코드
if __name__ == "__main__":
    parser = FileParser()
    
    print("FileParser 모듈 로드 완료")
    print(f"지원 형식: {parser.supported_formats}")